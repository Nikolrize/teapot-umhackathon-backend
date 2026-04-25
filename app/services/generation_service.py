import io
import json
import re
import unicodedata
import zipfile
from datetime import datetime

from app.services.glm_service import call_glm_session

# ── Latin-1 sanitiser (fpdf2 core fonts only support ISO-8859-1) ──────────────

_CHAR_MAP = str.maketrans({
    '—': '--',   # em dash  —
    '–': '-',    # en dash  –
    '‘': "'",    # left single quote  '
    '’': "'",    # right single quote '
    '“': '"',    # left double quote  "
    '”': '"',    # right double quote "
    '…': '...',  # ellipsis …
    '•': '-',    # bullet •
    '·': '-',    # middle dot ·
    '™': '(TM)', # trademark ™
    '®': '(R)',  # registered ®
    '©': '(c)',  # copyright ©
    '€': 'EUR',  # euro €
    '£': 'GBP',  # pound £
    '×': 'x',    # multiplication ×
    '≥': '>=',   # ≥
    '≤': '<=',   # ≤
})


def _safe(text: str) -> str:
    """Return a Latin-1-safe string for use with fpdf2 Helvetica/core fonts."""
    if not text:
        return ''
    text = text.translate(_CHAR_MAP)
    # Decompose accented chars (e.g. é -> e + combining accent), drop rest
    text = unicodedata.normalize('NFKD', text)
    return text.encode('latin-1', 'ignore').decode('latin-1')


# ── JSON schemas the LLM must follow ──────────────────────────────────────────

_PDF_PPT_SCHEMA = """
{
  "title": "Descriptive report title (max 10 words)",
  "subtitle": "Agent specialty — Business name",
  "executive_summary": "2–3 sentence professional summary of the most critical findings",
  "sections": [
    {
      "heading": "Section title",
      "narrative": "Professional prose paragraph with specific, context-grounded insights (3–5 sentences). Avoid generic statements.",
      "bullets": [
        "Specific, quantified insight or data point relevant to this business",
        "..."
      ],
      "table": {
        "caption": "Table title (only include when numeric or comparative data warrants it)",
        "headers": ["Column A", "Column B", "Column C"],
        "rows": [["value", "value", "value"]]
      }
    }
  ],
  "recommendations": [
    "Specific, prioritised, actionable recommendation with expected outcome — not generic advice"
  ],
  "conclusion": "Professional closing paragraph summarising key outlook and immediate next steps"
}"""

_CSV_SCHEMA = """
{
  "title": "Dataset title",
  "sheets": [
    {
      "name": "Sheet name (e.g. Summary, Projections, Risk Matrix)",
      "headers": ["Column 1", "Column 2", "Column 3"],
      "rows": [
        ["value", "value", "value"]
      ]
    }
  ]
}"""


# ── Prompt builder ─────────────────────────────────────────────────────────────

_DOC_LABELS = {
    "pdf": "PDF Business Report",
    "ppt": "PowerPoint Executive Presentation",
    "csv": "CSV Business Dataset",
}

_DOC_INSTRUCTIONS = {
    "pdf": (
        "Include 4–6 sections covering distinct aspects of the topic. "
        "Each section must contain a narrative paragraph, 3–5 bullet points, "
        "and a table wherever numeric or comparative data exists."
    ),
    "ppt": (
        "Include 4–6 sections. Each section becomes one slide. "
        "Keep slide headings under 6 words. "
        "Bullets must be concise (max 15 words each) — these appear on screen. "
        "Put detailed explanation in the narrative field (speaker notes)."
    ),
    "csv": (
        "Produce all relevant data tables a business analyst would need. "
        "Each logical grouping of data should be a separate sheet. "
        "Use realistic, context-grounded values — not placeholder data."
    ),
}


def build_prompts(
    session: dict,
    references: list,
    topic: str,
    doc_type: str,
) -> tuple[str, str]:
    """
    Returns (system_prompt, user_prompt).
    The system prompt establishes the agent's professional persona.
    The user prompt gives the structured output contract.
    """
    doc_label = _DOC_LABELS[doc_type]
    schema    = _CSV_SCHEMA if doc_type == "csv" else _PDF_PPT_SCHEMA
    doc_instr = _DOC_INSTRUCTIONS[doc_type]

    ref_block = ""
    if references:
        lines = "\n".join(f"    • {r['content']}" for r in references)
        ref_block = f"\nVerified reference data (treat as established fact — build on these, do not contradict):\n{lines}\n"

    effective_topic = topic or session["task"]

    system = (
        f"{session['requirements']}\n\n"
        f"You are now preparing a formal {doc_label} for {session['business_name']}, "
        f"a {session['business_type']} business. "
        f"This document will be presented to senior executives and board-level stakeholders. "
        f"Apply the highest standards of professional business writing:\n"
        f"  • Be specific and data-driven — reference the client's actual context\n"
        f"  • Be decisive — state clear findings, not hedged generalities\n"
        f"  • Be concise — every sentence must earn its place\n"
        f"  • Quantify wherever possible — percentages, timeframes, magnitudes\n"
        f"  • Avoid filler phrases such as 'it is important to note' or 'in conclusion'\n"
        f"  • Every recommendation must name a specific action, owner, and expected outcome"
    )

    user = (
        f"Generate a professional {doc_label} on the following topic.\n\n"
        f"TOPIC: {effective_topic}\n\n"
        f"CLIENT PROFILE:\n"
        f"  • Company:        {session['business_name']}\n"
        f"  • Industry:       {session['business_type']}\n"
        f"  • Business context: {session.get('business_context') or 'Not provided'}\n"
        f"  • Budget range:   {session.get('budget_min') or 'n/a'} – {session.get('budget_max') or 'n/a'}\n"
        f"  • Strategic goal: {session.get('goal') or 'Not provided'}\n"
        f"{ref_block}\n"
        f"DOCUMENT REQUIREMENTS:\n"
        f"  {doc_instr}\n\n"
        f"CRITICAL OUTPUT RULES:\n"
        f"  1. Return ONLY a raw JSON object. No markdown fences, no explanation, no preamble.\n"
        f"  2. Every figure and percentage must be grounded in the client context — not invented.\n"
        f"  3. Recommendations must be specific to {session['business_name']} — not generic templates.\n"
        f"  4. Omit the 'table' key from a section if no meaningful data table applies.\n\n"
        f"JSON SCHEMA (follow exactly):\n{schema}"
    )

    return system, user


# ── JSON extractor ─────────────────────────────────────────────────────────────

def _extract_json(text: str) -> dict:
    for attempt in (
        lambda: json.loads(text),
        lambda: json.loads(re.search(r'```(?:json)?\s*([\s\S]+?)\s*```', text).group(1)),
        lambda: json.loads(re.search(r'\{[\s\S]+\}', text).group(0)),
    ):
        try:
            return attempt()
        except Exception:
            continue
    raise ValueError(
        "The agent did not return valid structured JSON. "
        "Try again — if the problem persists, simplify the topic."
    )


# ── PDF builder ────────────────────────────────────────────────────────────────

def _build_pdf(data: dict) -> bytes:
    from fpdf import FPDF

    DARK  = (31,  56, 100)
    MID   = (68, 114, 196)
    GREY  = (245, 247, 250)
    TEXT  = (35,  35,  45)
    LIGHT_MID = (210, 220, 240)

    class _Doc(FPDF):
        def __init__(self, subtitle):
            super().__init__()
            self._sub = _safe(subtitle)

        def header(self):
            if self.page_no() == 1:
                return
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(160, 160, 160)
            self.cell(0, 8, self._sub, align="R")
            self.set_draw_color(210, 210, 210)
            self.line(10, self.get_y() + 8, 200, self.get_y() + 8)
            self.ln(10)

        def footer(self):
            if self.page_no() == 1:
                return
            self.set_y(-14)
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(160, 160, 160)
            self.cell(0, 8, f"Page {self.page_no()}", align="C")

    pdf = _Doc(data.get("subtitle", ""))
    pdf.set_auto_page_break(auto=True, margin=22)

    # ── Cover page ─────────────────────────────────────────────────────────────
    pdf.add_page()
    pdf.set_fill_color(*DARK)
    pdf.rect(0, 0, 210, 297, "F")

    # Accent stripe
    pdf.set_fill_color(*MID)
    pdf.rect(0, 130, 210, 4, "F")

    pdf.set_y(55)
    pdf.set_font("Helvetica", "B", 32)
    pdf.set_text_color(255, 255, 255)
    pdf.multi_cell(0, 14, _safe(data.get("title", "Business Report")), align="C")

    pdf.ln(6)
    pdf.set_font("Helvetica", "", 15)
    pdf.set_text_color(*LIGHT_MID)
    pdf.multi_cell(0, 8, _safe(data.get("subtitle", "")), align="C")

    pdf.set_y(250)
    pdf.set_font("Helvetica", "I", 10)
    pdf.set_text_color(130, 155, 200)
    pdf.cell(0, 8, f"Prepared  {datetime.now().strftime('%B %d, %Y')}", align="C")

    # ── Helpers for body pages ─────────────────────────────────────────────────
    def section_header(label: str):
        pdf.set_fill_color(*MID)
        pdf.set_font("Helvetica", "B", 12)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(0, 10, f"   {_safe(label)}", fill=True, ln=True)
        pdf.ln(3)

    def body_text(txt: str):
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(*TEXT)
        pdf.multi_cell(0, 6, _safe(txt))
        pdf.ln(2)

    def bullet_item(txt: str):
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(*TEXT)
        pdf.set_x(14)
        pdf.multi_cell(185, 6, f"-   {_safe(txt)}")

    def draw_table(tbl: dict):
        headers = tbl.get("headers", [])
        rows    = tbl.get("rows", [])
        caption = tbl.get("caption", "")
        if not headers:
            return
        if caption:
            pdf.set_font("Helvetica", "I", 9)
            pdf.set_text_color(110, 110, 110)
            pdf.cell(0, 5, _safe(caption), ln=True)
        col_w = min(185 / len(headers), 55)
        pdf.set_fill_color(*DARK)
        pdf.set_text_color(255, 255, 255)
        pdf.set_font("Helvetica", "B", 9)
        for h in headers:
            pdf.cell(col_w, 7, _safe(str(h))[:22], fill=True)
        pdf.ln()
        for i, row in enumerate(rows):
            pdf.set_fill_color(230, 237, 250) if i % 2 == 0 else pdf.set_fill_color(255, 255, 255)
            pdf.set_text_color(*TEXT)
            pdf.set_font("Helvetica", "", 9)
            for cell in row:
                pdf.cell(col_w, 6, _safe(str(cell))[:22], fill=True)
            pdf.ln()
        pdf.ln(3)

    # ── Executive summary ──────────────────────────────────────────────────────
    pdf.add_page()
    section_header("Executive Summary")
    pdf.set_fill_color(*GREY)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*TEXT)
    pdf.multi_cell(0, 7, _safe(data.get("executive_summary", "")), fill=True)
    pdf.ln(8)

    # ── Sections ───────────────────────────────────────────────────────────────
    for sec in data.get("sections", []):
        section_header(sec.get("heading", ""))
        if sec.get("narrative"):
            body_text(sec["narrative"])
        for b in sec.get("bullets", []):
            bullet_item(b)
        if sec.get("bullets"):
            pdf.ln(3)
        if sec.get("table"):
            draw_table(sec["table"])
        pdf.ln(4)

    # ── Recommendations ────────────────────────────────────────────────────────
    recs = data.get("recommendations", [])
    if recs:
        section_header("Recommendations")
        for i, rec in enumerate(recs, 1):
            pdf.set_x(12)
            pdf.set_font("Helvetica", "B", 10)
            pdf.set_text_color(*MID)
            pdf.cell(8, 7, f"{i}.")
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(*TEXT)
            pdf.multi_cell(180, 7, _safe(rec))
            pdf.ln(1)
        pdf.ln(4)

    # ── Conclusion ─────────────────────────────────────────────────────────────
    if data.get("conclusion"):
        section_header("Conclusion")
        pdf.set_font("Helvetica", "I", 10)
        pdf.set_text_color(*TEXT)
        pdf.multi_cell(0, 6, _safe(data["conclusion"]))

    return bytes(pdf.output())


# ── PPT builder ────────────────────────────────────────────────────────────────

def _build_pptx(data: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    DARK  = RGBColor(31,  56, 100)
    MID   = RGBColor(68, 114, 196)
    WHITE = RGBColor(255, 255, 255)
    LGREY = RGBColor(245, 247, 250)
    TEXT  = RGBColor(35,  35,  45)
    SUBTEXT = RGBColor(180, 200, 240)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide():
        return prs.slides.add_slide(blank)

    def bg(slide, color: RGBColor):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def bar(slide, top, height, color: RGBColor):
        tb = slide.shapes.add_textbox(Inches(0), top, Inches(13.33), height)
        tb.fill.solid()
        tb.fill.fore_color.rgb = color
        tb.line.fill.background()
        return tb

    def textbox(slide, text, x, y, w, h, size, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

    def bullet_frame(slide, bullets, x, y, w, h, size=14, color=TEXT):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•   {b}"
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(6)

    # ── Title slide ────────────────────────────────────────────────────────────
    sl = new_slide()
    bg(sl, DARK)
    bar(sl, Inches(3.4), Inches(0.06), MID)
    textbox(sl, data.get("title", "Business Report"),
            Inches(0.8), Inches(1.2), Inches(11.7), Inches(2),
            36, bold=True, color=WHITE)
    textbox(sl, data.get("subtitle", ""),
            Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.9),
            18, color=SUBTEXT)
    textbox(sl, datetime.now().strftime("%B %Y"),
            Inches(0.8), Inches(6.6), Inches(5), Inches(0.5),
            11, italic=True, color=RGBColor(130, 155, 200))

    # ── Executive summary ──────────────────────────────────────────────────────
    sl = new_slide()
    bg(sl, WHITE)
    bar(sl, Inches(0), Inches(1.25), DARK)
    textbox(sl, "Executive Summary",
            Inches(0.5), Inches(0.22), Inches(12), Inches(0.85),
            26, bold=True, color=WHITE)
    textbox(sl, data.get("executive_summary", ""),
            Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5),
            15, color=TEXT)

    # ── Section slides ─────────────────────────────────────────────────────────
    for sec in data.get("sections", []):
        sl = new_slide()
        bg(sl, WHITE)
        bar(sl, Inches(0), Inches(1.25), MID)
        textbox(sl, sec.get("heading", ""),
                Inches(0.5), Inches(0.22), Inches(12), Inches(0.85),
                22, bold=True, color=WHITE)

        bullets = sec.get("bullets", [])
        if bullets:
            bullet_frame(sl, bullets,
                         Inches(0.8), Inches(1.45), Inches(11.7), Inches(5.7))

        # Speaker notes from narrative
        if sec.get("narrative") and sl.has_notes_slide:
            sl.notes_slide.notes_text_frame.text = sec["narrative"]

    # ── Recommendations ────────────────────────────────────────────────────────
    recs = data.get("recommendations", [])
    if recs:
        sl = new_slide()
        bg(sl, WHITE)
        bar(sl, Inches(0), Inches(1.25), DARK)
        textbox(sl, "Recommendations",
                Inches(0.5), Inches(0.22), Inches(12), Inches(0.85),
                26, bold=True, color=WHITE)
        numbered = [f"{i}.   {r}" for i, r in enumerate(recs, 1)]
        bullet_frame(sl, numbered,
                     Inches(0.8), Inches(1.45), Inches(11.7), Inches(5.7),
                     size=14, color=TEXT)

    # ── Closing slide ──────────────────────────────────────────────────────────
    sl = new_slide()
    bg(sl, DARK)
    bar(sl, Inches(3.4), Inches(0.06), MID)
    textbox(sl, data.get("conclusion", ""),
            Inches(1), Inches(1.5), Inches(11.3), Inches(3.5),
            16, italic=True, color=SUBTEXT, align=PP_ALIGN.CENTER)
    textbox(sl, "Confidential — For Internal Use Only",
            Inches(0), Inches(6.8), Inches(13.33), Inches(0.5),
            9, italic=True, color=RGBColor(100, 130, 180), align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── CSV builder ────────────────────────────────────────────────────────────────

def _build_csv(data: dict) -> tuple[bytes, str]:
    """Returns (bytes, filename). Zips if multiple sheets are returned."""
    import pandas as pd

    sheets = data.get("sheets", [])
    if not sheets:
        return b"", "empty.csv"

    if len(sheets) == 1:
        s = sheets[0]
        df = pd.DataFrame(s.get("rows", []), columns=s.get("headers", []))
        return df.to_csv(index=False).encode("utf-8"), f"{s.get('name', 'data')}.csv"

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for s in sheets:
            df = pd.DataFrame(s.get("rows", []), columns=s.get("headers", []))
            name = s.get("name", "Sheet").replace(" ", "_") + ".csv"
            zf.writestr(name, df.to_csv(index=False))
    return buf.getvalue(), f"{data.get('title', 'report').replace(' ', '_')[:40]}.zip"


# ── Main entry point ───────────────────────────────────────────────────────────

def generate_document(
    session: dict,
    references: list,
    topic: str,
    doc_type: str,
    model_info: dict | None = None,
) -> tuple[bytes, str, str]:
    """
    Runs the two-step pipeline:
      1. LLM call with professional prompt → structured JSON
      2. Document builder → file bytes

    Returns (file_bytes, filename, mime_type).
    """
    system, user = build_prompts(session, references, topic, doc_type)

    raw, tokens_used = call_glm_session(
        4096,
        system,
        [{"role": "user", "content": user}],
        temperature=0.65,
        top_p=0.90,
        api_key=model_info["api_key"]        if model_info else None,
        model_name=model_info["model_name"]  if model_info else None,
        model_provider=model_info["model_provider"] if model_info else None,
    )

    data = _extract_json(raw)

    safe_title = re.sub(r"[^\w\s-]", "", data.get("title", "report"))[:40].strip().replace(" ", "_")

    if doc_type == "pdf":
        return _build_pdf(data), f"{safe_title}.pdf", "application/pdf", tokens_used
    if doc_type == "ppt":
        return (
            _build_pptx(data),
            f"{safe_title}.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            tokens_used,
        )
    file_bytes, filename = _build_csv(data)
    mime = "application/zip" if filename.endswith(".zip") else "text/csv"
    return file_bytes, filename, mime, tokens_used
